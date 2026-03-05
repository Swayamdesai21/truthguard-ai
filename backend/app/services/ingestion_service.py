"""
Ingestion Service — Multi-format document chunking and indexing.

Handles:
- Text extraction from PDF (PyPDF2), DOCX (python-docx), PPTX (python-pptx)
- Plain text and Markdown files
- Sentence-aware text chunking
- Building BM25 + TF-IDF indices and storing in SessionStore
"""
from __future__ import annotations
import io
import re
from typing import List, Optional
import PyPDF2

from app.services.vector_store import get_or_create_session, set_session, SessionStore
from app.services.embedding_service import build_indices
from app.core.config import CHUNK_SIZE, CHUNK_OVERLAP


# ── Text Extractors ──────────────────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract all text from a PDF file given its raw bytes."""
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract all text from a DOCX file."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        return "\n\n".join(paragraphs)
    except Exception as e:
        print(f"[Ingestion] DOCX extraction error: {e}")
        return ""


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract all text from a PPTX (PowerPoint) file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {slide_num}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
                # Also extract from tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_parts.append(row_text)
            if len(slide_parts) > 1:  # more than just the slide header
                slides_text.append("\n".join(slide_parts))
        return "\n\n".join(slides_text)
    except Exception as e:
        print(f"[Ingestion] PPTX extraction error: {e}")
        return ""


# Map of supported extensions → extractor functions
EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".doc": extract_text_from_docx,  # best-effort; true .doc may not parse
    ".pptx": extract_text_from_pptx,
    ".ppt": extract_text_from_pptx,   # best-effort
}

# Extensions treated as plain text
PLAIN_TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".csv", ".json", ".xml", ".html", ".htm", ".log", ".rst"}


def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """
    Auto-detect file format and extract text.
    Returns extracted text or empty string on failure.
    """
    ext = _get_extension(filename)

    if ext in EXTRACTORS:
        return EXTRACTORS[ext](file_bytes)
    elif ext in PLAIN_TEXT_EXTENSIONS:
        return _decode_text(file_bytes)
    else:
        # Try to decode as text (fallback)
        return _decode_text(file_bytes)


def _get_extension(filename: str) -> str:
    """Get lowercase file extension."""
    import os
    _, ext = os.path.splitext(filename.lower())
    return ext


def _decode_text(content: bytes) -> str:
    """Try to decode bytes as text with multiple encodings."""
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            return content.decode(encoding)
        except (UnicodeDecodeError, ValueError):
            continue
    return ""


# ── Chunking ─────────────────────────────────────────────────────────────────

def sentence_aware_chunk(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """
    Split text into overlapping chunks, breaking at sentence boundaries.
    """
    text = re.sub(r"\s+", " ", text).strip()
    sentence_endings = re.compile(r"(?<=[.!?])\s+")
    sentences = sentence_endings.split(text)

    chunks: List[str] = []
    current_chunk: List[str] = []
    current_len = 0

    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue

        sentence_len = len(sentence)

        if current_len + sentence_len + 1 > chunk_size and current_chunk:
            chunk_text = " ".join(current_chunk)
            chunks.append(chunk_text)

            overlap_text = chunk_text[-overlap:] if len(chunk_text) > overlap else chunk_text
            overlap_sentences = sentence_endings.split(overlap_text)
            current_chunk = [s.strip() for s in overlap_sentences if s.strip()]
            current_len = sum(len(s) for s in current_chunk)

        current_chunk.append(sentence)
        current_len += sentence_len + 1

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return [c for c in chunks if len(c.strip()) > 20]


# ── Ingest Functions ─────────────────────────────────────────────────────────

def ingest_text(session_id: str, text: str, source_name: str = "document") -> int:
    """Ingest raw text into the session store. Returns chunk count."""
    chunks = sentence_aware_chunk(text)
    if not chunks:
        return 0

    store = get_or_create_session(session_id)

    start_idx = len(store.chunks)
    store.chunks.extend(chunks)
    store.metadata.extend([
        {"source": source_name, "chunk_index": start_idx + i}
        for i in range(len(chunks))
    ])

    bm25, vectorizer, tfidf_matrix = build_indices(store.chunks)
    store.bm25_index = bm25
    store.tfidf_vectorizer = vectorizer
    store.tfidf_matrix = tfidf_matrix

    set_session(session_id, store)
    return len(chunks)


def ingest_pdf(session_id: str, pdf_bytes: bytes, source_name: str = "document.pdf") -> int:
    """Extract text from PDF and ingest."""
    text = extract_text_from_pdf(pdf_bytes)
    if not text.strip():
        return 0
    return ingest_text(session_id, text, source_name)


def ingest_file(session_id: str, file_bytes: bytes, filename: str) -> int:
    """
    Auto-detect format, extract text, and ingest into session.
    Supports: PDF, DOCX, PPTX, TXT, MD, CSV, and more.
    Returns chunk count.
    """
    text = extract_text_from_file(file_bytes, filename)
    if not text.strip():
        return 0
    return ingest_text(session_id, text, source_name=filename)
